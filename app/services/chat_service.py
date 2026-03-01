from __future__ import annotations

import base64
import io
import logging
import mimetypes
import uuid
from typing import AsyncGenerator

from sqlalchemy import select

logger = logging.getLogger(__name__)
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.models.message import MediaAttachment, Message
from app.models.thread import Thread
from app.services.llm.base import TokenUsage
from app.services.llm.router import get_provider
from app.services.llm.status import provider_status_tracker
from app.services.pricing_service import MediaCounts, compute_cost, get_current_price
from app.storage.base import StorageBackend

MAX_HISTORY_MESSAGES = 50

async def load_thread_history(
    db: AsyncSession, thread: Thread, include_media: bool = True,
) -> list[dict]:
    """Load the most recent messages from a thread and format them for LLM consumption."""
    stmt = (
        select(Message)
        .where(Message.thread_id == thread.id)
        .order_by(Message.created_at.desc())
        .limit(MAX_HISTORY_MESSAGES)
        .options(selectinload(Message.attachments))
    )
    result = await db.execute(stmt)
    messages = list(reversed(result.scalars().all()))

    formatted: list[dict] = []
    for msg in messages:
        if not include_media or not msg.attachments:
            formatted.append({"role": msg.role, "content": msg.content or ""})
            continue

        has_images = any(a.media_type == "image" for a in msg.attachments)
        has_text_files = any(a.text_content for a in msg.attachments)

        if has_images or has_text_files:
            content_parts: list[dict] = []
            if msg.content:
                content_parts.append({"type": "text", "text": msg.content})
            for att in msg.attachments:
                if att.media_type == "image":
                    content_parts.append({
                        "type": "image_url",
                        "image_url": {"url": f"attachment://{att.id}"},
                    })
                elif att.text_content:
                    content_parts.append({
                        "type": "text",
                        "text": f"[File: {att.mime_type}]\n{att.text_content}",
                    })
            formatted.append({"role": msg.role, "content": content_parts or msg.content})
        else:
            formatted.append({"role": msg.role, "content": msg.content or ""})
    return formatted


def _flatten_text_only_parts(messages: list[dict]) -> list[dict]:
    """If a message content is a list but contains only text parts, merge into a single string."""
    result: list[dict] = []
    for msg in messages:
        if isinstance(msg["content"], list):
            has_non_text = any(p.get("type") != "text" for p in msg["content"])
            if not has_non_text:
                merged = "\n\n".join(p.get("text", "") for p in msg["content"])
                result.append({"role": msg["role"], "content": merged})
                continue
        result.append(msg)
    return result


async def build_llm_messages(
    db: AsyncSession,
    thread: Thread,
    user_content: str | list[dict],
    storage: StorageBackend,
) -> list[dict]:
    """Build the full message list from thread history.

    The current user message must already be persisted (via save_user_message)
    before calling this, so it is included in the history query.
    """
    history = await load_thread_history(db, thread)
    history = await _resolve_attachments(db, history, storage)
    history = _flatten_text_only_parts(history)
    logger.debug("LLM message count: %d, last msg role: %s, content len: %d",
                 len(history), history[-1]["role"] if history else "?",
                 len(str(history[-1].get("content", ""))) if history else 0)
    return history


async def _resolve_attachments(
    db: AsyncSession, messages: list[dict], storage: StorageBackend,
) -> list[dict]:
    """Replace attachment:// URLs with base64 data URLs for vision models."""
    resolved: list[dict] = []
    for msg in messages:
        if isinstance(msg["content"], list):
            parts: list[dict] = []
            for part in msg["content"]:
                if (
                    part.get("type") == "image_url"
                    and part.get("image_url", {}).get("url", "").startswith("attachment://")
                ):
                    att_id = part["image_url"]["url"].removeprefix("attachment://")
                    result = await db.execute(
                        select(MediaAttachment).where(MediaAttachment.id == uuid.UUID(att_id))
                    )
                    att = result.scalar_one_or_none()
                    if att:
                        path = await storage.get_path(att.file_path)
                        if path.is_file():
                            with open(path, "rb") as f:
                                data = base64.b64encode(f.read()).decode()
                            mime = att.mime_type or "image/jpeg"
                            parts.append({
                                "type": "image_url",
                                "image_url": {"url": f"data:{mime};base64,{data}"},
                            })
                            continue
                        else:
                            logger.warning("Attachment file missing from storage: %s", att.file_path)
                    else:
                        logger.warning("Attachment record not found: %s", att_id)
                    continue
                parts.append(part)
            resolved.append({"role": msg["role"], "content": parts})
        else:
            resolved.append(msg)
    return resolved


AUDIO_BYTES_PER_SECOND = 16_000  # ~128 kbps
VIDEO_BYTES_PER_SECOND = 1_000_000  # ~8 Mbps rough estimate


def _build_media_counts(attachments: list[MediaAttachment] | None) -> MediaCounts:
    if not attachments:
        return MediaCounts()
    images = 0
    audio_sec = 0.0
    video_sec = 0.0
    for att in attachments:
        if att.media_type == "image":
            images += 1
        elif att.media_type == "audio":
            audio_sec += att.file_size / AUDIO_BYTES_PER_SECOND
        elif att.media_type == "video":
            video_sec += att.file_size / VIDEO_BYTES_PER_SECOND
    return MediaCounts(
        image_count=images,
        audio_seconds=audio_sec,
        video_seconds=video_sec,
    )


async def stream_llm_response(
    db: AsyncSession,
    thread: Thread,
    storage: StorageBackend,
    attachments: list[MediaAttachment] | None = None,
) -> AsyncGenerator[dict, None]:
    """Stream LLM response chunks as dicts suitable for SSE serialization.

    The current user message must already be persisted before calling this.

    Event types:
    - {"event": "chunk", "data": "<text>"}
    - {"event": "done", "data": {"content": "...", "model": "...", "message_id": "..."}}
    - {"event": "error", "data": "<error message>"}
    """
    try:
        provider = get_provider(thread.llm_name)
        messages = await build_llm_messages(db, thread, "", storage)

        full_response: list[str] = []
        usage: TokenUsage | None = None
        async for chunk in provider.stream_completion(messages, model=thread.llm_name):
            if isinstance(chunk, TokenUsage):
                usage = chunk
            else:
                full_response.append(chunk)
                yield {"event": "chunk", "data": chunk}

        content = "".join(full_response)
        provider_status_tracker.record_success(provider.provider_name())

        cost_usd = None
        if usage:
            pricing = await get_current_price(db, thread.llm_name)
            if pricing:
                media = _build_media_counts(attachments)
                cost_usd = compute_cost(usage, pricing, media)

        assistant_msg = Message(
            thread_id=thread.id,
            role="assistant",
            content=content,
            model=thread.llm_name,
            prompt_tokens=usage.prompt_tokens if usage else None,
            completion_tokens=usage.completion_tokens if usage else None,
            token_count=usage.total_tokens if usage else None,
            web_search_calls=usage.web_search_calls if usage else None,
            tool_calls=usage.tool_calls if usage else None,
            cost_usd=cost_usd,
        )
        db.add(assistant_msg)
        await db.commit()
        await db.refresh(assistant_msg)

        done_data: dict = {
            "content": content,
            "model": thread.llm_name,
            "message_id": str(assistant_msg.id),
        }
        if usage:
            done_data["usage"] = {
                "prompt_tokens": usage.prompt_tokens,
                "completion_tokens": usage.completion_tokens,
                "total_tokens": usage.total_tokens,
            }
        if cost_usd is not None:
            done_data["cost_usd"] = float(cost_usd)

        yield {"event": "done", "data": done_data}
    except Exception as exc:
        try:
            p = get_provider(thread.llm_name)
            provider_status_tracker.record_failure(p.provider_name(), str(exc))
        except ValueError:
            pass
        yield {"event": "error", "data": str(exc)}


async def save_user_message(
    db: AsyncSession,
    thread: Thread,
    content: str,
    attachments: list[MediaAttachment] | None = None,
) -> Message:
    """Persist a user message and its attachments."""
    msg = Message(thread_id=thread.id, role="user", content=content)
    db.add(msg)
    await db.flush()

    if attachments:
        for att in attachments:
            att.message_id = msg.id
            db.add(att)

    await db.commit()
    await db.refresh(msg)
    return msg


MAX_EXTRACTED_CHARS = 100_000

_TEXT_MIME_TYPES = {
    "application/json", "application/xml", "application/javascript",
    "application/x-yaml", "application/csv", "application/sql",
    "application/x-sh", "application/x-python", "application/rtf",
}

_DOCX_TYPES = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
}
_XLSX_TYPES = {
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.ms-excel",
}
_PPTX_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
}


def _extract_text(data: bytes, content_type: str, filename: str) -> str | None:
    """Extract readable text from a file. Returns None if not applicable."""
    logger.info("_extract_text: filename=%s, content_type=%s, size=%d bytes",
                filename, content_type, len(data))
    lower = filename.lower()

    if content_type == "application/pdf" or lower.endswith(".pdf"):
        result = _extract_pdf(data, filename)
        logger.info("PDF extraction result: %s", f"{len(result)} chars" if result else "None")
        return result

    if content_type in _DOCX_TYPES or lower.endswith(".docx"):
        return _extract_docx(data, filename)

    if content_type in _XLSX_TYPES or lower.endswith((".xlsx", ".xls")):
        return _extract_xlsx(data, filename)

    if content_type in _PPTX_TYPES or lower.endswith(".pptx"):
        return _extract_pptx(data, filename)

    if lower.endswith(".csv") or content_type == "text/csv":
        return _extract_plain(data)

    if content_type.startswith("text/") or content_type in _TEXT_MIME_TYPES:
        return _extract_plain(data)

    if lower.endswith((".txt", ".md", ".log", ".ini", ".cfg", ".toml",
                       ".yml", ".yaml", ".json", ".xml", ".html", ".htm",
                       ".css", ".js", ".ts", ".py", ".java", ".c", ".cpp",
                       ".h", ".cs", ".go", ".rs", ".rb", ".php", ".sh",
                       ".bat", ".ps1", ".sql", ".r", ".m", ".swift")):
        return _extract_plain(data)

    return None


def _extract_plain(data: bytes) -> str | None:
    try:
        text = data.decode("utf-8", errors="replace").strip()
        return text[:MAX_EXTRACTED_CHARS] if text else None
    except Exception:
        return None


def _extract_pdf(data: bytes, filename: str) -> str | None:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        logger.info("PDF has %d pages", len(reader.pages))
        pages = [p.extract_text() for p in reader.pages if p.extract_text()]
        combined = "\n\n".join(pages).strip()
        return combined[:MAX_EXTRACTED_CHARS] if combined else None
    except Exception:
        logger.exception("PDF text extraction failed for %s", filename)
        return None


def _extract_docx(data: bytes, filename: str) -> str | None:
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))
        combined = "\n".join(paragraphs).strip()
        return combined[:MAX_EXTRACTED_CHARS] if combined else None
    except Exception:
        logger.warning("DOCX text extraction failed for %s", filename)
        return None


def _extract_xlsx(data: bytes, filename: str) -> str | None:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sections: list[str] = []
        for sheet in wb.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                header = f"[Sheet: {sheet.title}]"
                sections.append(header + "\n" + "\n".join(rows))
        wb.close()
        combined = "\n\n".join(sections).strip()
        return combined[:MAX_EXTRACTED_CHARS] if combined else None
    except Exception:
        logger.warning("XLSX text extraction failed for %s", filename)
        return None


def _extract_pptx(data: bytes, filename: str) -> str | None:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        combined = "\n\n".join(slides).strip()
        return combined[:MAX_EXTRACTED_CHARS] if combined else None
    except Exception:
        logger.warning("PPTX text extraction failed for %s", filename)
        return None


async def process_uploaded_files(
    files_data: list[tuple[str, bytes, str]],
    storage: StorageBackend,
    thread_id: uuid.UUID,
) -> list[MediaAttachment]:
    """Save uploaded files and create MediaAttachment records (not yet linked to a message)."""
    attachments: list[MediaAttachment] = []
    subdir = str(thread_id)

    for original_name, data, content_type in files_data:
        ext = mimetypes.guess_extension(content_type) or ""
        filename = f"{uuid.uuid4().hex}{ext}"

        key = await storage.save(data, filename, subdir)

        media_type = "image" if content_type.startswith("image/") else \
                     "video" if content_type.startswith("video/") else \
                     "audio" if content_type.startswith("audio/") else "file"

        thumbnail_key = None
        if media_type == "image":
            try:
                thumbnail_key = await storage.save_thumbnail(data, filename, subdir)
            except Exception:
                pass

        text_content = None
        if media_type == "file":
            text_content = _extract_text(data, content_type, original_name)
            logger.info("File %s: media_type=%s, text_content=%s",
                        original_name, media_type,
                        f"{len(text_content)} chars" if text_content else "None")

        att = MediaAttachment(
            media_type=media_type,
            file_path=key,
            original_filename=original_name,
            mime_type=content_type,
            file_size=len(data),
            thumbnail_path=thumbnail_key,
            text_content=text_content,
        )
        attachments.append(att)

    return attachments
